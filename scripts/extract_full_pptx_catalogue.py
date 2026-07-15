from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import re, os, json, csv

pptx_path = r'C:\Users\91779\Downloads\catalouge full.pptx'
out_dir = r'C:\Users\91779\TresorCouture\Tresor-Couture\inventory-from-pptx'
os.makedirs(out_dir, exist_ok=True)
img_dir = os.path.join(out_dir, 'images-full')
os.makedirs(img_dir, exist_ok=True)
for f in os.listdir(img_dir):
    os.remove(os.path.join(img_dir, f))

public_img_dir = r'C:\Users\91779\TresorCouture\Tresor-Couture\public\products\pptx-full'
os.makedirs(public_img_dir, exist_ok=True)
for f in os.listdir(public_img_dir):
    os.remove(os.path.join(public_img_dir, f))

prs = Presentation(pptx_path)

def clean_text(t):
    return re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', ' ', t)

def get_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            t = shape.text.strip()
            if t:
                texts.append(t)
    return '\n'.join(texts)

def parse_code(text):
    # explicit CODE first (letters optional)
    m = re.search(r'code\s*[-:\u2013\u2014]?\s*([A-Z]{0,3}\s?\d{3,5})', text, re.IGNORECASE)
    if m:
        return re.sub(r'\s+', '', m.group(1)).upper()
    # fallback: any code-like token with at least two letters, not preceded by price/metric keywords
    for m in re.finditer(r'(?<!\w)([A-Z]{2,3}\s?\d{3,5})(?!\w)', text):
        cand = m.group(1)
        before = text[:m.start()].strip()[-15:].upper()
        if any(before.endswith(p) for p in ['BP','SP','BUNDLE','B','M','STOCK','OG','METER','METERS']):
            continue
        if re.search(r'(?:BP|SP|BUNDLE|B|M|STOCK|OG|METER|METERS)\s*[-:]?\s*$', before, re.I):
            continue
        # don't treat bare "M 480" / "B 1299" as codes
        if re.search(r'\b[mMbB]\s*$', before):
            continue
        return re.sub(r'\s+', '', cand).upper()
    return None

def parse_number(text, patterns):
    for pat in patterns:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            return m.group(1).strip()
    return None

inventory = []
seen_codes = {}
for idx, slide in enumerate(prs.slides, 1):
    text = clean_text(get_slide_text(slide))
    print(f'\n--- Slide {idx} ---')
    print(text)

    code = parse_code(text)
    # Make code unique by appending slide suffix if duplicate
    if code:
        base_code = code
        if base_code in seen_codes:
            code = f'{base_code}-S{idx}'
        seen_codes[base_code] = True
    else:
        code = f'NO-CODE-S{idx}'

    bp = parse_number(text, [
        r'(?:bp|buying price)\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
    ])

    # Per-meter / per-unit "M" price. Must not be part of "meter(s)".
    sp_per_meter = parse_number(text, [
        r'(?:sp|selling price)\s*(?:per meter|/meter|/ meter|meter|m)?\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'sp/meter\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'sp\s*(\d+(?:\.\d+)?)',
    ])
    if not sp_per_meter:
        # Standalone "M 380", "M 480", or "METER 30" used as a per-unit price
        m_price = re.search(r'(?<![a-zA-Z])m(?:eter)?\s*(\d+(?:\.\d+)?)(?=\s*(?:\n|$|bp|sp|b\s+\d|og|bundle|stock|/|\\))', text, re.IGNORECASE)
        if m_price:
            sp_per_meter = m_price.group(1)

    bundle = parse_number(text, [
        r'(?:sp\s+only\s+)?bundle\s*(?:price)?\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'only\s+bundle\s*(\d+(?:\.\d+)?)',
        r'(?<![a-zA-Z])b\s*(\d+(?:\.\d+)?)(?=\s*(?:\n|$|sold|only|\(|\s+\d))',
        r'bundle\s+of\s+\d+(?:\.\d+)?\s*(?:meter|meters|m)\s+(?:at\s+)?(?:Rs\.?\s*)?(\d+(?:\.\d+)?)',
    ])

    og = parse_number(text, [
        r'(?<![a-zA-Z])og\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
    ])

    stock = parse_number(text, [
        r'stock\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)\s*(?:meter|m)?',
        r'(\d+(?:\.\d+)?)\s*(?:meter|meters?|m)\b',
    ])

    colour = parse_number(text, [
        r'colou?rs?\s*[-:\u2013\u2014]?\s*([A-Za-z\-/\s]+?)(?:\n|$|stock|bp|sp|bundle|meter|og|available|\d)',
    ])

    colour_text = colour
    notes = []
    if colour_text:
        raw = colour_text
        colour_text = re.sub(r'\(.*$', '', colour_text).strip()
        colour_text = re.sub(r'\d+\s+available', '', colour_text, flags=re.I).strip()
        colour_text = colour_text.strip('-').strip()
        # discard non-colour tokens captured by the loose regex
        if colour_text.lower() in {'available', 'sp', 's', 'colours', 'colors', 'colour', 'color'}:
            colour_text = ''
        if re.search(r'\d+\s*available', raw, re.I) and not colour_text:
            notes.append('Multiple colours available')

    # Determine if metered (lace) or unit (patch)
    has_meters = bool(stock) and float(stock) > 0
    is_lace = has_meters or 'meter' in text.lower() or 'meters' in text.lower()

    # Bundle logic
    explicit_bundle = bundle
    bundle_price = bundle
    if is_lace and not bundle_price and sp_per_meter and stock:
        try:
            bundle_price = str(int(round(float(sp_per_meter) * float(stock))))
        except Exception:
            pass

    # Determine listing price (what customer pays)
    # Explicit bundle price always wins: it is the fixed price for all meters/units in stock.
    if explicit_bundle:
        listing_price = float(explicit_bundle)
        unit_type = 'bundle'
    elif sp_per_meter:
        listing_price = float(sp_per_meter)
        unit_type = 'per meter' if is_lace else 'unit'
    elif bp and not is_lace:
        # Non-metered / bundle-only items: BP is the bundle selling price.
        listing_price = float(bp)
        unit_type = 'bundle'
    elif og:
        listing_price = float(og)
        unit_type = 'unit'
    else:
        listing_price = 0
        unit_type = 'unit'

    # MRP = 30% extra so that price is 30% off MRP
    # price = 0.7 * mrp  =>  mrp = price / 0.7
    mrp = round(listing_price / 0.7, 2) if listing_price else 0

    image_files = []
    for s_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                ext = img.ext or 'png'
                fname = f'slide{idx:02d}_img{s_idx}_{code}.{ext}'
                fpath = os.path.join(img_dir, fname)
                with open(fpath, 'wb') as f:
                    f.write(img.blob)
                image_files.append(fpath)
                public_path = os.path.join(public_img_dir, fname)
                with open(public_path, 'wb') as f:
                    f.write(img.blob)
                print(f'  Saved image: {fname}')
            except Exception as e:
                print(f'  Image error: {e}')

    if code.startswith('NO-CODE'):
        notes.append('Code not found in slide text; using generated code')
    # clean accidental price-marker prefixes from code
    if code:
        code = re.sub(r'^(M|B|SP|BP|OG)(\d)', r'\2', code)
    if not listing_price:
        notes.append('No selling price found')

    row = {
        'slide': idx,
        'code': code,
        'bp': bp,
        'sp_per_meter': sp_per_meter,
        'bundle_price': bundle_price,
        'og': og,
        'stock_meters': stock,
        'colour_from_text': colour_text if colour_text else '',
        'unit_type': unit_type,
        'is_lace': is_lace,
        'listing_price': listing_price,
        'mrp': mrp,
        'notes': '; '.join(notes),
        'raw_text': text,
        'image_files': ';'.join(os.path.basename(p) for p in image_files),
    }
    inventory.append(row)

# CSV
csv_path = os.path.join(out_dir, 'inventory_full.csv')
with open(csv_path, 'w', newline='', encoding='utf-8') as f:
    writer = csv.DictWriter(f, fieldnames=inventory[0].keys())
    writer.writeheader()
    writer.writerows(inventory)

# JSON
json_path = os.path.join(out_dir, 'inventory_full.json')
with open(json_path, 'w', encoding='utf-8') as f:
    json.dump(inventory, f, indent=2, ensure_ascii=False)

# Seed JSON for Firestore
fabrics = []
for r in inventory:
    if not r['listing_price']:
        continue
    colour = r['colour_from_text'] or ''
    stock = int(float(r['stock_meters'] or 1))
    category_word = 'lace' if r['is_lace'] else 'patch / embellishment'
    price_display = f"₹{int(r['listing_price'])}"

    if r['unit_type'] == 'bundle':
        if r['is_lace']:
            bundle_desc = f"Sold as a complete bundle of {stock} meter{'s' if stock != 1 else ''} at {price_display}."
        else:
            bundle_desc = f"Sold as a complete bundle at {price_display}."
        description = (
            f"Elegant {category_word} from TRESOR. "
            f"{bundle_desc} "
            f"Perfect for premium couture and bridal projects."
        )
    elif r['unit_type'] == 'per meter':
        description = (
            f"Elegant {category_word} from TRESOR. "
            f"Priced at {price_display} per meter; {stock} meter{'s' if stock != 1 else ''} currently available. "
            f"Ideal for premium couture and bridal projects."
        )
    else:
        description = (
            f"Elegant {category_word} from TRESOR. "
            f"Available at {price_display}. "
            f"Perfect for premium couture and bridal projects."
        )

    if colour.strip():
        description += f" Colour: {colour.strip()}."

    if r['notes']:
        description += f" Note: {r['notes']}"

    fabrics.append({
        'id': r['code'],
        'brand': 'TRESOR',
        'name': r['code'],
        'productCode': r['code'],
        'description': description,
        'price': r['listing_price'],
        'mrp': r['mrp'],
        'photo': f"/products/pptx-full/{r['image_files'].split(';')[0]}" if r['image_files'] else '',
        'image': '',
        'category': 'Laces',
        'masterCategory': 'Laces',
        'subCategory': 'Trim & Edging',
        'materialType': 'Mixed',
        'tags': ['New In'],
        'colors': [{'name': colour.strip(), 'hex': '#cccccc'}] if colour.strip() else [],
        'stock': stock,
        'unitType': r['unit_type'],
        'costPrice': None,
        'sellingPricePerMeter': float(r['sp_per_meter']) if r['sp_per_meter'] else None,
        'bundleSizeMeters': stock if r['unit_type'] == 'bundle' and stock else None,
        'bundlePrice': float(r['bundle_price']) if r['bundle_price'] else None,
    })
seed_path = os.path.join(out_dir, 'inventory_full_seed.json')
with open(seed_path, 'w', encoding='utf-8') as f:
    json.dump(fabrics, f, indent=2, ensure_ascii=False)

print(f'\nDone. CSV: {csv_path}, JSON: {json_path}, Seed: {seed_path}, Images: {img_dir}')
print(f'Total products: {len(fabrics)}')
