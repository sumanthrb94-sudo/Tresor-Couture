from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import numpy as np
import re, os, json, csv, math

pptx_path = r'C:\Users\91779\Downloads\catalouge.pptx'
out_dir = r'C:\Users\91779\TresorCouture\Tresor-Couture\inventory-from-pptx'
os.makedirs(out_dir, exist_ok=True)
img_dir = os.path.join(out_dir, 'images')
os.makedirs(img_dir, exist_ok=True)
# Copy images to public/products/pptx/ so they are served by the live site
public_img_dir = r'C:\Users\91779\TresorCouture\Tresor-Couture\public\products\pptx'
os.makedirs(public_img_dir, exist_ok=True)
for f in os.listdir(public_img_dir):
    os.remove(os.path.join(public_img_dir, f))
# clean previous images for a fresh run
for f in os.listdir(img_dir):
    os.remove(os.path.join(img_dir, f))

prs = Presentation(pptx_path)

colors_rgb = {
    'Black': (0,0,0), 'White': (255,255,255), 'Red': (255,0,0), 'Maroon': (128,0,0),
    'Wine': (114,47,55), 'Pink': (255,192,203), 'Rose': (255,0,128),
    'Purple': (128,0,128), 'Lavender': (230,230,250), 'Violet': (238,130,238),
    'Blue': (0,0,255), 'Navy': (0,0,128), 'Sky Blue': (135,206,235), 'Cyan': (0,255,255),
    'Teal': (0,128,128), 'Green': (0,128,0), 'Lime': (0,255,0), 'Olive': (128,128,0),
    'Yellow': (255,255,0), 'Gold': (255,215,0), 'Orange': (255,165,0), 'Peach': (255,218,185),
    'Brown': (165,42,42), 'Beige': (245,245,220), 'Cream': (255,253,208), 'Ivory': (255,255,240),
    'Grey': (128,128,128), 'Silver': (192,192,192), 'Charcoal': (54,69,79),
    'Magenta': (255,0,255), 'Turquoise': (64,224,208),
}

def clean_text(t):
    return re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', ' ', t)

def get_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            t = shape.text.strip()
            if t:
                texts.append(t)
    return '\n'.join(texts)

def parse_code(text):
    m = re.search(r'code\s*[-:\u2013\u2014]?\s*([A-Z]{1,3}\s?\d{3,5})', text, re.IGNORECASE)
    if m:
        return re.sub(r'\s+', '', m.group(1)).upper()
    for m in re.finditer(r'(?<!\w)([A-Z]{1,3}\s?\d{3,5})(?!\w)', text):
        cand = m.group(1)
        before = text[:m.start()].strip()[-15:].upper()
        if any(before.endswith(p) for p in ['BP','SP','BUNDLE','B','M','STOCK']):
            continue
        if re.search(r'(?:BP|SP|BUNDLE|B|M|STOCK)\s*[-:]?\s*$', before, re.I):
            continue
        return re.sub(r'\s+', '', cand).upper()
    return None

def parse_number(text, patterns):
    for pat in patterns:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            return m.group(1).strip()
    return None

def nearest_color(rgb):
    r,g,b = [int(x) for x in rgb]
    best = None
    best_d = float('inf')
    for name, (cr,cg,cb) in colors_rgb.items():
        d = (r-cr)**2 + (g-cg)**2 + (b-cb)**2
        if d < best_d:
            best_d = d
            best = name
    return best

def dominant_colors(img_path, n=3):
    try:
        img = Image.open(img_path).convert('RGB')
        w,h = img.size
        left, top = int(w*0.2), int(h*0.2)
        right, bottom = int(w*0.8), int(h*0.8)
        crop = img.crop((left, top, right, bottom))
        small = crop.resize((60,60))
        arr = np.array(small).reshape(-1,3)
        mask = ~(
            ((arr[:,0]>245)&(arr[:,1]>245)&(arr[:,2]>245)) |
            ((arr[:,0]<25)&(arr[:,1]<25)&(arr[:,2]<25)) |
            ((abs(arr[:,0]-arr[:,1])<15)&(abs(arr[:,1]-arr[:,2])<15)&(arr[:,0]>120))
        )
        filtered = arr[mask]
        if len(filtered) < 50:
            filtered = arr
        bins = {}
        for px in filtered:
            key = (px[0]//10*10, px[1]//10*10, px[2]//10*10)
            bins[key] = bins.get(key, 0) + 1
        top = sorted(bins.items(), key=lambda x: -x[1])[:n]
        return [nearest_color(tuple(c)) for c,_ in top]
    except Exception as e:
        return [f'error:{e}']

inventory = []
for idx, slide in enumerate(prs.slides, 1):
    if idx > 15:
        break
    text = clean_text(get_slide_text(slide))
    print(f'\n--- Slide {idx} ---')
    print(text)

    code = parse_code(text)
    bp = parse_number(text, [
        r'(?:bp|buying price)\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
    ])
    sp_per_meter = parse_number(text, [
        r'(?:sp|selling price)\s*(?:per meter|/meter|/ meter|meter|m)?\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'sp/meter\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'meter\s*(\d+(?:\.\d+)?)',
        r'\bm\s*(\d+(?:\.\d+)?)',
    ])
    stock = parse_number(text, [
        r'stock\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)\s*(?:meter|m)?',
    ])
    colour = parse_number(text, [
        r'colou?r\s*[-:\u2013\u2014]?\s*([A-Za-z\-/\s]+?)(?:\n|$|stock|bp|sp|bundle|meter)',
    ])
    bundle = parse_number(text, [
        r'bundle\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'bundle\s+of\s+\d+\s+meters\s+(\d+(?:\.\d+)?)',
    ])

    colour_text = colour
    notes = []
    if colour_text:
        raw = colour_text
        colour_text = re.sub(r'\(.*$', '', colour_text).strip()
        colour_text = re.sub(r'\d+\s+available', '', colour_text, flags=re.I).strip()
        colour_text = colour_text.strip('-').strip()
        if re.search(r'\d+\s*available', raw, re.I) and not colour_text:
            notes.append('Multiple colours available')

    # Treat explicit "SP bundle" mentions as the bundle price if no general bundle price found
    sp_bundle = parse_number(text, [
        r'sp\s*bundle\s*[-:\u2013\u2014]?\s*(\d+(?:\.\d+)?)',
        r'sp\s*only\s*bundle\s*price\s*(\d+(?:\.\d+)?)',
    ])
    if not bundle and sp_bundle:
        bundle = sp_bundle

    image_files = []
    for s_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                ext = img.ext or 'png'
                fname = f'slide{idx:02d}_img{s_idx}_{code or "no_code"}.{ext}'
                fpath = os.path.join(img_dir, fname)
                with open(fpath, 'wb') as f:
                    f.write(img.blob)
                image_files.append(fpath)
                public_path = os.path.join(public_img_dir, fname)
                with open(public_path, 'wb') as f:
                    f.write(img.blob)
                print(f'  Saved image: {fname}')
            except Exception as e:
                print(f'  Image error: {e}')

    image_colors = []
    if not colour_text and image_files:
        image_colors = dominant_colors(image_files[0])
        print(f'  Dominant colors from image: {image_colors}')

    if not code:
        notes.append('Code not found in slide text')

    row = {
        'slide': idx,
        'code': code,
        'bp': bp,
        'sp_per_meter': sp_per_meter,
        'stock_meters': stock,
        'colour_from_text': colour_text if colour_text else '',
        'colour_from_image': ', '.join(image_colors) if image_colors else '',
        'bundle_price': bundle,
        'notes': '; '.join(notes),
        'raw_text': text,
        'image_files': ';'.join(os.path.basename(p) for p in image_files),
    }
    inventory.append(row)

csv_path = os.path.join(out_dir, 'inventory_first15.csv')
with open(csv_path, 'w', newline='', encoding='utf-8') as f:
    writer = csv.DictWriter(f, fieldnames=inventory[0].keys())
    writer.writeheader()
    writer.writerows(inventory)

json_path = os.path.join(out_dir, 'inventory_first15.json')
with open(json_path, 'w', encoding='utf-8') as f:
    json.dump(inventory, f, indent=2, ensure_ascii=False)

fabrics = []
for r in inventory:
    colour = r['colour_from_text'] or (r['colour_from_image'].split(',')[0] if r['colour_from_image'] else '')
    # If no per-meter SP exists, the product is sold as a bundle; use bundle price as the listing price.
    listing_price = float(r['sp_per_meter'] or r['bundle_price'] or r['bp'] or 0)
    fabrics.append({
        'id': r['code'] or f'SLIDE{r["slide"]:02d}',
        'brand': 'TRESOR',
        'name': r['code'] or f'Product Slide {r["slide"]}',
        'description': f'Extracted from catalogue slide {r["slide"]}. {r["notes"]}',
        'price': listing_price,
        'mrp': listing_price,
        'photo': f'/products/pptx/{r["image_files"].split(";")[0]}' if r['image_files'] else '',
        'image': '',
        'category': 'Laces',
        'masterCategory': 'Laces',
        'subCategory': 'Trim & Edging',
        'materialType': 'Mixed',
        'origin': 'India',
        'tags': ['New In'],
        'colors': [{'name': colour.strip(), 'hex': '#cccccc'}] if colour.strip() else [],
        'stock': int(float(r['stock_meters'] or 0)),
    })
seed_path = os.path.join(out_dir, 'inventory_first15_seed.json')
with open(seed_path, 'w', encoding='utf-8') as f:
    json.dump(fabrics, f, indent=2, ensure_ascii=False)

print(f'\nDone. CSV: {csv_path}, JSON: {json_path}, Seed: {seed_path}, Images: {img_dir}')
